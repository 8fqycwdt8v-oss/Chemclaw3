"""Turn a document's bytes into text, structurally and offline. The one parsing implementation.

These parsers began life inside `chemclaw.agent.attachments`, serving a chemist handing the agent a
file. They live here now because reading a PDF is an *ingest* concern that an upload happens to
use, not the other way round — and because `chemclaw.ingest` may not import `chemclaw.agent`
(`tests/test_layering.py`), so a mounted-share crawler could otherwise only have got at them by
growing a second copy. One implementation, two callers, no drift.

**The honesty rule that governs every parser here.** Extraction is *structural*, never heuristic:
page, slide, sheet and cell boundaries come from each format's own document model, and a file the
library cannot open is refused rather than salvaged. A PDF "parsed" by scraping text-like bytes
yields confident nonsense a chemist cannot distinguish from a real reading.

That rule is why a **scanned** PDF is refused by name (`ScannedDocumentError`) rather than returned
as an empty document. Silence must never read as "the file was blank"; the failure mode being
avoided is a chemist concluding a CoA had nothing in it. On a decade-old file share this is not a
corner case — it is a population, which is why the refusal has its own exception type: the share
sync counts scans separately and reports how much of the corpus is invisible.

No format is accepted whose reading needs a network service (D-089).
"""

import csv
import io
import logging
import zipfile
from collections.abc import Callable

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pydantic import BaseModel
from pypdf import PdfReader

from chemclaw.core.config import settings
from chemclaw.ingest.documents.formats import EXTENSIONS, content_type_for

logger = logging.getLogger(__name__)


class DocumentParseError(ValueError):
    """A document that cannot be read, with a message naming what is supported."""


class ScannedDocumentError(DocumentParseError):
    """A PDF with no text layer at all — a scan or an image-only export.

    Its own type because it is the one refusal that is about the *document* rather than about this
    system's format support: the file is a perfectly good PDF that simply contains no text. A share
    sync counts these apart from unsupported formats so an operator can see how much of the corpus
    would need OCR, instead of one undifferentiated "skipped" number.
    """


class ParsedDocument(BaseModel):
    """One document's extracted text, with the content type it was read as."""

    content_type: str
    text: str
    # Row count for a tabular document, so a caller can say "42 runs" without re-parsing.
    rows: int = 0


def _refuse_a_bomb(name: str, raw: bytes) -> None:
    """Refuse an OOXML container whose parts expand past the configured ceiling.

    `.docx`/`.xlsx`/`.pptx` are zip archives, so every size limit upstream of here — a share's
    `max_file_bytes`, an upload's `attachment_max_bytes` — bounds only the *compressed* bytes. A
    110 KB workbook holding 31 MB of sheet XML is a 282× ratio, measured; at the share's 50 MB
    default that is ~14 GB of XML and multiple GB of Python strings, and the worker is OOM-killed
    with no counter, no log line and no report.

    Read from the central directory, so it costs no decompression. **The residual is stated:** a
    hand-crafted archive can understate `file_size`, and this check believes it. That is a bound on
    the realistic case — a real generator writes true sizes — not a defence against a crafted one,
    which needs a streaming limit at every read.

    Raises:
        DocumentParseError: The declared expansion exceeds `document_max_expanded_bytes`.
    """
    ceiling = settings.document_max_expanded_bytes
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as container:
            expanded = sum(item.file_size for item in container.infolist())
    except zipfile.BadZipFile as exc:
        raise DocumentParseError(f"could not read {name}: {exc}") from exc
    if expanded > ceiling:
        raise DocumentParseError(
            f"{name} expands to {expanded} bytes from {len(raw)} on disk, past the "
            f"{ceiling}-byte limit. A document this large compressed this well is a data export or "
            "a malformed file rather than a document; extracting the relevant sheet will work."
        )


def _parse_text(raw: bytes) -> tuple[str, int]:
    """Decode a text document verbatim — nothing is summarized or dropped at ingest."""
    return raw.decode("utf-8", errors="replace"), 0


def _parse_csv(raw: bytes) -> tuple[str, int]:
    """Render a delimited table as aligned text, preserving every cell.

    Rendered rather than handed over as raw CSV because the agent reads prose far more reliably
    than it reads quoting rules, and because a mangled quote in a raw paste can silently shift a
    whole column — a wrong number a chemist would have no way to spot.
    """
    text = raw.decode("utf-8", errors="replace")
    dialect_sample = text[:4096]
    try:
        dialect = csv.Sniffer().sniff(dialect_sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel  # a single-column or unusual file is still readable as plain rows
    rows = list(csv.reader(io.StringIO(text), dialect))
    if not rows:
        return "", 0
    header, *body = rows
    lines = [" | ".join(header), "-" * 40]
    lines += [" | ".join(cell for cell in row) for row in body]
    return "\n".join(lines), len(body)


def _parse_pdf(raw: bytes) -> tuple[str, int]:
    """Extract a PDF's text layer page by page; refuse a scan rather than return nothing.

    Pages are labelled and kept in order because a chemist citing "the table on page 3" needs the
    page to survive ingest — an unlabelled concatenation loses the only coordinate the source
    document offers.

    A PDF where **no page** yields text is a scan (or an image-only export): `pypdf` reports
    success and returns nothing, which would present to the agent as an *empty document* and to the
    chemist as "there was nothing in it". Refused by name instead, since that is a true statement
    about what this system can read and the empty text is not.

    The test is "did any page produce text at all", deliberately not a minimum length. A one-line
    CoA is a legitimate upload, and any threshold tuned to document size would refuse it — the
    thing that distinguishes a scan is that it yields *zero* characters, not few.
    """
    reader = PdfReader(io.BytesIO(raw))
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    if not any(pages):
        raise ScannedDocumentError(
            f"no text could be extracted from any of this PDF's {len(pages)} page(s), so it is a "
            "scan or an image-only export. Reading it needs OCR, which is not built — a text-based "
            "PDF, or the relevant text pasted directly, will work."
        )
    # Page labels come from the original numbering, so a page that is itself a scan drops out
    # without renumbering the ones after it — a citation to "page 3" must still land on page 3.
    return "\n\n".join(
        f"[page {number}]\n{text}" for number, text in enumerate(pages, 1) if text
    ), len(pages)


def _parse_pptx(raw: bytes) -> tuple[str, int]:
    """Extract a deck's text slide by slide, including tables and speaker notes.

    Notes are included because a project deck's reasoning frequently lives there rather than on the
    slide, and dropping them would silently discard the most informative half of the file.
    """
    deck = Presentation(io.BytesIO(raw))
    blocks: list[str] = []
    slides = list(deck.slides)
    for number, slide in enumerate(slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                parts += [" | ".join(cell.text for cell in row.cells) for row in shape.table.rows]
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"(speaker notes) {notes}")
        if parts:
            blocks.append(f"[slide {number}]\n" + "\n".join(parts))
    return "\n\n".join(blocks), len(slides)


def _parse_docx(raw: bytes) -> tuple[str, int]:
    """Extract a Word document's paragraphs and tables in document order.

    Tables are rendered with the same `|` separator `_parse_csv` uses, so a table reads identically
    however it reached the system — one representation for the agent to learn, not three.
    """
    document = Document(io.BytesIO(raw))
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    rows = 0
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
            rows += 1
    return "\n".join(parts), rows


def _parse_xlsx(raw: bytes) -> tuple[str, int]:
    """Extract a workbook sheet by sheet as delimited rows.

    `data_only=True` reads the *cached values* of formula cells rather than the formulas: a chemist
    attaching a yield sheet means the yields, and `=B2/C2*100` is not an answer. A workbook saved
    without cached values yields empty cells there, which is visible rather than wrong.
    """
    book = load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
    try:
        blocks: list[str] = []
        rows = 0
        for sheet in book.worksheets:
            lines = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    lines.append(" | ".join("" if c is None else str(c) for c in row))
                    rows += 1
            if lines:
                blocks.append(f"[sheet {sheet.title}]\n" + "\n".join(lines))
        return "\n\n".join(blocks), rows
    finally:
        # read_only workbooks hold an open zip handle; leaking it would exhaust file descriptors
        # over a long-lived pod's worth of uploads.
        book.close()


# The closed allowlist. A content type absent here is refused with a message, never guessed at.
_PARSERS: dict[str, Callable[[bytes], tuple[str, int]]] = {
    "text/markdown": _parse_text,
    "text/plain": _parse_text,
    "text/csv": _parse_csv,
    "text/tab-separated-values": _parse_csv,
    "application/pdf": _parse_pdf,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _parse_pptx,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _parse_docx,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _parse_xlsx,
}

# The zip-container formats. Their size limits upstream bound compressed bytes only, so these are
# the three that need the expansion check before a parser is handed the archive.
_ZIP_CONTAINERS = frozenset(
    {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
)

# The two halves of the allowlist must name the same formats: `formats.EXTENSIONS` decides what a
# crawl even opens, `_PARSERS` decides what can be read. A format in one and not the other is
# either a file type that silently never matches or a parser nothing can reach — both invisible at
# runtime, both caught here at import.
_UNPARSEABLE = set(EXTENSIONS.values()) - set(_PARSERS)
_UNREACHABLE = set(_PARSERS) - set(EXTENSIONS.values())
if _UNPARSEABLE or _UNREACHABLE:  # pragma: no cover - a wiring error, not a runtime state
    raise ImportError(
        "the format allowlist and the parser table disagree; "
        f"declared but unparseable: {sorted(_UNPARSEABLE)}; "
        f"parseable but undeclared: {sorted(_UNREACHABLE)}"
    )


def parse_document(name: str, raw: bytes, declared_type: str | None = None) -> ParsedDocument:
    """Extract one document's text, or refuse it with a message naming the supported formats.

    Args:
        name: The file name; its extension resolves the format when no type is declared.
        raw: The document's bytes.
        declared_type: A content type from the transport, if the caller has one.

    Returns:
        The extracted text with the content type it was read as.

    Raises:
        ScannedDocumentError: A PDF with no text layer at all.
        DocumentParseError: An unsupported format, or a file that could not be read.
    """
    content_type = content_type_for(name, declared_type)
    parser = _PARSERS.get(content_type)
    if parser is None:
        raise DocumentParseError(
            f"{name} ({content_type}) is not a supported format. Supported: "
            f"{', '.join(sorted(EXTENSIONS))}. Spectra and image formats need OCR/vision "
            "ingestion, which is not built — exporting the relevant text or table is the "
            "reliable path today."
        )
    if content_type in _ZIP_CONTAINERS:
        _refuse_a_bomb(name, raw)
    try:
        text, rows = parser(raw)
    except DocumentParseError:
        # Already precise — a refusal the parser named itself, `ScannedDocumentError` included.
        raise
    except Exception as exc:
        # **One net, at the boundary, around the whole parse.** Each parser used to guard only its
        # *constructor*, which is the one call that is not where these libraries do their work:
        # `openpyxl` in `read_only` mode parses the sheet inside `iter_rows`, and `python-pptx`
        # loads slide parts lazily. A truncated `sheet1.xml` from an interrupted network copy
        # therefore raised `ElementTree.ParseError` — a `SyntaxError`, caught by nothing — and an
        # unbalanced quote in an instrument export raised `csv.Error` from the reader, likewise
        # outside every guard. On the share those escaped the sync's reject-and-continue net, failed
        # the activity, and since the crawl keeps no cross-run cursor, **every later run restarted
        # and hit the same file**: one malformed document stopped the whole corpus indexing.
        #
        # Broad on purpose, and only here: `raw` is untrusted bytes from a share or an upload, and
        # every library below is a third-party parser over them. "This file could not be read" is
        # the honest statement about any failure in that region, and it is the statement the callers
        # already handle — one counted refusal rather than a dead job.
        raise DocumentParseError(f"could not read {name}: {exc}") from exc
    return ParsedDocument(content_type=content_type, text=text, rows=rows)
