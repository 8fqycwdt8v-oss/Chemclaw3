"""PDF/PPTX/DOCX/XLSX ingest (D-089), and the refusal that survived the scope change.

These formats were originally refused outright: a PDF "parsed" by scraping text-like bytes yields
confident nonsense a chemist cannot distinguish from a real reading. The scope decision reversed
the refusal, not the reasoning — so what these tests pin is that extraction is *structural* (each
format read through its own document model, with pages/slides/sheets preserved) and that the one
case real extraction cannot fix, a **scanned** PDF, is still refused by name rather than returned
as an empty document.

Every fixture is **built here by the format's own writer**, never a checked-in blob. A committed
binary would make these tests assert something about a file someone once produced; building the
document in the test means the assertion is about our parsing and stays honest if a library
changes.
"""

import io

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from agents.attachments import AttachmentError, content_type_for, parse_attachment


def _docx_bytes(paragraphs: list[str], table: list[list[str]] | None = None) -> bytes:
    """A .docx carrying the given paragraphs and an optional table."""
    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    if table:
        word_table = document.add_table(rows=len(table), cols=len(table[0]))
        for word_row, row in zip(word_table.rows, table, strict=True):
            for cell, value in zip(word_row.cells, row, strict=True):
                cell.text = value
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(slides: list[tuple[str, str]], notes: str = "") -> bytes:
    """A .pptx with one title+body slide per entry, and optional notes on the first."""
    deck = Presentation()
    for title, body in slides:
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
        box.text_frame.text = body
        if notes and slide is deck.slides[0]:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes(sheets: dict[str, list[list[object]]]) -> bytes:
    """An .xlsx with one named worksheet per entry."""
    book = Workbook()
    book.remove(book.active)
    for name, rows in sheets.items():
        sheet = book.create_sheet(title=name)
        for row in rows:
            sheet.append(row)
    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()


def _blank_pdf_bytes(pages: int = 2) -> bytes:
    """A PDF with real pages and no text layer — what a scan looks like to a parser."""
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=595, height=842)
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def _text_pdf_bytes(pages: list[str]) -> bytes:
    """A PDF carrying one line of real, extractable text per page.

    Assembled by hand rather than with a renderer: `pypdf` writes PDFs but cannot *typeset*, and
    pulling in reportlab purely to generate fixtures would add a dependency the shipped code never
    uses. This is the smallest structure that produces a genuine text layer — catalog, page tree,
    one content stream per page holding a `BT … Tj ET` text object, and a correct xref table.
    """
    objects: list[bytes] = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        "<</Type/Pages/Kids[{}]/Count {}>>".format(
            " ".join(f"{3 + 2 * i} 0 R" for i in range(len(pages))), len(pages)
        ).encode(),
    ]
    font_id = 3 + 2 * len(pages)
    for index, text in enumerate(pages):
        content_id = 4 + 2 * index
        objects.append(
            f"<</Type/Page/Parent 2 0 R/MediaBox[0 0 595 842]/Contents {content_id} 0 R"
            f"/Resources<</Font<</F1 {font_id} 0 R>>>>>>".encode()
        )
        stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
        objects.append(b"<</Length %d>>\nstream\n%s\nendstream" % (len(stream), stream))
    objects.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % number + body + b"\nendobj\n"
    xref_at = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objects) + 1)
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += b"trailer\n<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objects) + 1,
        xref_at,
    )
    return bytes(out)


# --- PDF ----------------------------------------------------------------------------------------


def test_a_pdf_with_a_text_layer_is_extracted_page_by_page() -> None:
    """The primary case: a real report reaches the agent as its actual words, not an approximation.

    This is what makes the format safe to accept at all — the original refusal was about *guessing*
    at bytes, and extraction through the document's own text objects is not guessing.
    """
    raw = _text_pdf_bytes(["Yield 84 percent", "Impurity below 0.5 percent"])
    attachment = parse_attachment("coa.pdf", raw)
    assert attachment.rows == 2  # page count
    assert "Yield 84 percent" in attachment.text
    assert "Impurity below 0.5 percent" in attachment.text


def test_pdf_page_numbers_are_preserved_so_a_citation_still_resolves() -> None:
    """A citation to page 3 has to still resolve to page 3 after ingest."""
    text = parse_attachment("r.pdf", _text_pdf_bytes(["alpha", "beta", "gamma"])).text
    assert "[page 1]" in text
    assert "[page 3]" in text
    assert text.index("alpha") < text.index("beta") < text.index("gamma")


def test_a_mixed_pdf_keeps_its_text_pages_and_does_not_invent_the_blank_ones() -> None:
    """A report with a scanned figure page must not have that page silently renumbered away.

    The blank page contributes no text — but it must not shift the labels of the pages that follow
    it, or a citation to "page 3" would land on page 2's content.
    """
    reader_input = _text_pdf_bytes(["first page text", "", "third page text"])
    text = parse_attachment("mixed.pdf", reader_input).text
    assert "[page 1]" in text
    assert "[page 3]" in text
    assert "[page 2]" not in text  # it had nothing to show, and nothing was fabricated for it


def test_a_scanned_pdf_is_refused_rather_than_read_as_an_empty_document() -> None:
    """The one case proper extraction cannot fix, and the reason the original refusal existed.

    `pypdf` opens an image-only PDF happily and returns nothing, which would reach the agent as a
    document with no content — and the chemist as "there was nothing in your CoA". The refusal is
    the only truthful answer, so it names the cause instead of the file being silently blank.
    """
    with pytest.raises(AttachmentError) as excinfo:
        parse_attachment("scan.pdf", _blank_pdf_bytes())
    message = str(excinfo.value)
    assert "scan" in message
    assert "OCR" in message
    # Names the evidence it judged on, so a chemist can tell this from a size or format refusal.
    assert "2 page(s)" in message


def test_a_short_pdf_is_accepted_because_the_scan_test_is_zero_text_not_a_length() -> None:
    """A one-line CoA is a legitimate document; any length threshold would refuse it.

    The distinguishing property of a scan is that it yields *no* characters, not few — so this
    pins that the refusal above cannot drift into a size check.
    """
    attachment = parse_attachment("short.pdf", _text_pdf_bytes(["ok"]))
    assert "ok" in attachment.text


def test_a_corrupt_pdf_is_refused_with_a_message_not_a_traceback() -> None:
    """A malformed upload is a user error; it must not surface as an unhandled library exception."""
    with pytest.raises(AttachmentError, match="could not read the PDF"):
        parse_attachment("broken.pdf", b"%PDF-1.4 this is not actually a pdf")


# --- PPTX ---------------------------------------------------------------------------------------


def test_a_deck_is_read_slide_by_slide_with_boundaries_preserved() -> None:
    """Slide numbers survive ingest: "the third slide" has to still mean something afterwards."""
    raw = _pptx_bytes([("Route A", "Pd(OAc)2, 80 C"), ("Route B", "NiCl2, 100 C")])
    attachment = parse_attachment("routes.pptx", raw)
    assert attachment.rows == 2  # slide count
    assert "[slide 1]" in attachment.text
    assert "[slide 2]" in attachment.text
    assert "Pd(OAc)2, 80 C" in attachment.text
    assert attachment.text.index("Route A") < attachment.text.index("Route B")


def test_speaker_notes_are_extracted_because_the_reasoning_often_lives_there() -> None:
    """A deck's "why" is usually in the notes; dropping them discards the informative half."""
    raw = _pptx_bytes([("Screen", "12 conditions")], notes="Ligand screen failed on the chloride.")
    text = parse_attachment("screen.pptx", raw).text
    assert "Ligand screen failed on the chloride." in text
    assert "(speaker notes)" in text


# --- DOCX ---------------------------------------------------------------------------------------


def test_a_word_document_yields_its_paragraphs_and_its_tables() -> None:
    """Both bodies of content, not just the prose — a protocol's numbers live in its table."""
    raw = _docx_bytes(
        ["Procedure for the Suzuki coupling.", "Charge the vessel under nitrogen."],
        table=[["reagent", "equiv"], ["boronic acid", "1.2"]],
    )
    attachment = parse_attachment("sop.docx", raw)
    assert "Charge the vessel under nitrogen." in attachment.text
    assert "boronic acid | 1.2" in attachment.text
    assert attachment.rows == 2


def test_a_docx_table_reads_the_same_as_a_csv_table() -> None:
    """One table representation for the agent to learn, whichever format delivered it."""
    docx_text = parse_attachment("t.docx", _docx_bytes([], table=[["a", "b"], ["1", "2"]])).text
    csv_text = parse_attachment("t.csv", b"a,b\n1,2\n").text
    assert "a | b" in docx_text
    assert "a | b" in csv_text


# --- XLSX ---------------------------------------------------------------------------------------


def test_a_workbook_is_read_sheet_by_sheet_with_names_kept() -> None:
    """Sheet names carry meaning in a lab workbook ("batch 3"), so they survive ingest."""
    raw = _xlsx_bytes({"yields": [["run", "yield"], [1, 84], [2, 91]], "notes": [["ok"]]})
    attachment = parse_attachment("runs.xlsx", raw)
    assert "[sheet yields]" in attachment.text
    assert "[sheet notes]" in attachment.text
    assert "1 | 84" in attachment.text
    assert attachment.rows == 4


def test_empty_rows_are_dropped_but_empty_cells_are_kept_as_gaps() -> None:
    """A blank row is spreadsheet padding; a blank *cell* is a missing measurement — not the same.

    Collapsing an empty cell would shift every value after it into the wrong column, which is the
    silent-wrong-number failure the CSV parser's docstring warns about.
    """
    raw = _xlsx_bytes({"s": [["a", "b", "c"], [], [1, None, 3]]})
    text = parse_attachment("gaps.xlsx", raw).text
    assert "1 |  | 3" in text
    assert parse_attachment("gaps.xlsx", raw).rows == 2


# --- the allowlist ------------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("name", "expected"),
    [
        ("report.pdf", "application/pdf"),
        ("deck.PPTX", "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        ("sop.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        ("runs.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    ],
)
def test_the_extension_resolves_when_a_browser_declares_nothing_useful(
    name: str, expected: str
) -> None:
    """Uploads routinely arrive as `application/octet-stream`; the extension has to carry it."""
    assert content_type_for(name, "application/octet-stream") == expected


def test_the_allowlist_is_still_closed_and_refuses_by_name() -> None:
    """Widening the scope did not turn the allowlist into a guess-what-this-is parser.

    The refusal now advertises the extensions rather than the MIME types: a chemist knows their
    file is a `.jdx`, and a list of `application/vnd.openxmlformats-…` strings tells them nothing.
    """
    with pytest.raises(AttachmentError) as excinfo:
        parse_attachment("spectrum.jdx", b"##TITLE=NMR\n")
    message = str(excinfo.value)
    assert ".pdf" in message
    assert ".xlsx" in message
    assert "OCR" in message


def test_a_binary_upload_never_reaches_a_text_parser() -> None:
    """The old failure mode: bytes that are not text being decoded into plausible-looking noise."""
    with pytest.raises(AttachmentError):
        parse_attachment("image.png", b"\x89PNG\r\n\x1a\n" + b"\x00" * 64)
