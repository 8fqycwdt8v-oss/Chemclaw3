"""The document fixtures every parsing test builds on — one writer per format.

Every fixture is **built by the format's own writer**, never a checked-in blob. A committed binary
would make a test assert something about a file someone once produced; building the document here
means the assertion is about our parsing, and stays honest if a library changes.

Shared rather than per-file because two suites need them now — the format tests
(`test_document_formats.py`) and the mounted-share tests (`test_document_share.py`) — and a second
copy would be a second set of documents to keep in step with the parsers.
"""

import io

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter


def _docx_bytes(paragraphs: list[str], table: list[list[str]] | None = None) -> bytes:
    """A .docx carrying the given paragraphs and an optional table."""
    document = Document()
    for text in paragraphs:
        document.add_paragraph(text)
    if table:
        word_table = document.add_table(rows=len(table), cols=len(table[0]))
        for word_row, row in zip(word_table.rows, table, strict=True):
            for cell, value in zip(word_row.cells, row, strict=True):
                cell.text = value
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(slides: list[tuple[str, str]], notes: str = "") -> bytes:
    """A .pptx with one title+body slide per entry, and optional notes on the first."""
    deck = Presentation()
    for title, body in slides:
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(2))
        box.text_frame.text = body
        if notes and slide is deck.slides[0]:
            slide.notes_slide.notes_text_frame.text = notes
    buffer = io.BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes(sheets: dict[str, list[list[object]]]) -> bytes:
    """An .xlsx with one named worksheet per entry."""
    book = Workbook()
    book.remove(book.active)
    for name, rows in sheets.items():
        sheet = book.create_sheet(title=name)
        for row in rows:
            sheet.append(row)
    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()


def _blank_pdf_bytes(pages: int = 2) -> bytes:
    """A PDF with real pages and no text layer — what a scan looks like to a parser."""
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=595, height=842)
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def _text_pdf_bytes(pages: list[str]) -> bytes:
    """A PDF carrying one line of real, extractable text per page.

    Assembled by hand rather than with a renderer: `pypdf` writes PDFs but cannot *typeset*, and
    pulling in reportlab purely to generate fixtures would add a dependency the shipped code never
    uses. This is the smallest structure that produces a genuine text layer — catalog, page tree,
    one content stream per page holding a `BT … Tj ET` text object, and a correct xref table.
    """
    objects: list[bytes] = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        "<</Type/Pages/Kids[{}]/Count {}>>".format(
            " ".join(f"{3 + 2 * i} 0 R" for i in range(len(pages))), len(pages)
        ).encode(),
    ]
    font_id = 3 + 2 * len(pages)
    for index, text in enumerate(pages):
        content_id = 4 + 2 * index
        objects.append(
            f"<</Type/Page/Parent 2 0 R/MediaBox[0 0 595 842]/Contents {content_id} 0 R"
            f"/Resources<</Font<</F1 {font_id} 0 R>>>>>>".encode()
        )
        stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
        objects.append(b"<</Length %d>>\nstream\n%s\nendstream" % (len(stream), stream))
    objects.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n" % number + body + b"\nendobj\n"
    xref_at = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objects) + 1)
    for offset in offsets:
        out += b"%010d 00000 n \n" % offset
    out += b"trailer\n<</Size %d/Root 1 0 R>>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objects) + 1,
        xref_at,
    )
    return bytes(out)


def _with_truncated_member(container: bytes, suffix: str) -> bytes:
    """A structurally valid OOXML file whose named part is cut in half.

    What an interrupted network copy leaves on a share: the zip directory is intact, so the file
    opens, and the damage is only found when the part is actually parsed. That distinction is the
    whole point — `openpyxl` in `read_only` mode and `python-pptx` both parse lazily, well after the
    constructor a guard would naturally sit on.
    """
    import zipfile

    source = zipfile.ZipFile(io.BytesIO(container))
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as out:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename.endswith(suffix):
                data = data[: len(data) // 2]
            out.writestr(item, data)
    return buffer.getvalue()


def _unbalanced_quote_csv_bytes(field_chars: int = 200_000) -> bytes:
    """A CSV whose stray quote swallows the rest of the file into one oversized field.

    An instrument export with one unescaped `"` does this. Everything after it becomes a single
    field, and `csv` refuses past its 131072-character limit — from the reader, not the sniffer.
    """
    return b'name,value\n"' + b"x" * field_chars


def _highly_compressible_xlsx_bytes(rows: int = 20_000) -> bytes:
    """A workbook that is small on disk and large in memory — the ratio, not the size, is the point.

    Every size limit upstream of the parser bounds *compressed* bytes, and OOXML is a zip. Repeated
    identical rows compress to almost nothing, so this passes `max_file_bytes` and
    `attachment_max_bytes` comfortably while expanding by a couple of orders of magnitude.

    **The row count is chosen against the two properties its caller asserts, not for size.** The
    test refuses the workbook against a `document_max_expanded_bytes` of 1 MB and measures
    `ratio > 20`. Building 200,000 rows through openpyxl cost 5.5 s of the suite for margin nobody
    used — measured: 35.7 MB expanded, ratio 22.8. 20,000 rows costs 0.43 s and gives 3.5 MB
    expanded at ratio 21.8: still 3.5× the cap and clear of the ratio floor. 10,000 would also pass
    (1.75 MB, ratio 21.0) and is not taken, because 5,000 measures 19.9 — one step below the
    assertion — and a fixture one step from failing is a fixture that fails on an openpyxl release.
    """
    book = Workbook()
    sheet = book.active
    for _ in range(rows):
        sheet.append(["aaaaaaaaaaaaaaaaaaaaaaaaaaaaaa", "bbbbbbbbbbbbbbbbbbbbbbbbbbbbbb"])
    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()
